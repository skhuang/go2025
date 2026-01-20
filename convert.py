import os
import sys
import argparse
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.util import Inches, Pt

def convert_pdf_to_pptx(pdf_path, pptx_path, editable=False):
    """
    Converts a PDF file to a PowerPoint presentation.
    If editable=False: Converts each page to an image and embeds it.
    If editable=True: Attempts to reconstruct the layout with editable text and specific images.
    """
    if not os.path.exists(pdf_path):
        print(f"Error: PDF file not found at {pdf_path}")
        return

    print(f"Converting {pdf_path} to {pptx_path} (Editable: {editable})...")

    try:
        doc = fitz.open(pdf_path)
    except Exception as e:
        print(f"Error opening PDF: {e}")
        return

    prs = Presentation()
    
    # Use a blank slide layout
    blank_slide_layout = prs.slide_layouts[6]

    for i in range(len(doc)):
        print(f"Processing page {i + 1}/{len(doc)}...")
        page = doc.load_page(i)
        
        slide = prs.slides.add_slide(blank_slide_layout)
        
        # Determine slide dimensions
        slide_width = prs.slide_width
        slide_height = prs.slide_height

        if editable:
            # Reconstruct layout
            blocks = page.get_text("dict")["blocks"]
            
            # Pdf units are usually Points (1/72 inch). PPTX uses EMU or Inches/Pt.
            # python-pptx Pt takes points.
            
            # We need to scale PDF coordinates to Slide coordinates if page size differs.
            # Usually we just assume Point to Point mapping is 1:1, but let's check page size.
            page_rect = page.rect
            pdf_width = page_rect.width
            pdf_height = page_rect.height
            
            # Scale factor if we want to fit the PDF page into the PPTX slide 10x7.5 inches
            # Or we can resize the slide to match the PDF.
            # Resizing slide is better for fidelity.
            prs.slide_width = int(pdf_width * 12700) # 1 pt = 12700 EMUs
            prs.slide_height = int(pdf_height * 12700)

            for block in blocks:
                bbox = block["bbox"]
                b_x0, b_y0, b_x1, b_y1 = bbox
                b_width = b_x1 - b_x0
                b_height = b_y1 - b_y0
                
                if block["type"] == 0: # Text
                    # Combined text for the block
                    text_content = ""
                    for line in block["lines"]:
                        for span in line["spans"]:
                            text_content += span["text"] + " "
                        text_content += "\n"
                    
                    # Add text box
                    # box coordinates in points, need separate unit if not using Pt
                    txBox = slide.shapes.add_textbox(Pt(b_x0), Pt(b_y0), Pt(b_width), Pt(b_height))
                    tf = txBox.text_frame
                    tf.text = text_content.strip()
                    
                    # Try to apply some simple formatting from the first span (simplified)
                    # Ideally we would iterate spans and add runs, but blocks are often grouped.
                    try:
                        first_span = block["lines"][0]["spans"][0]
                        tf.paragraphs[0].font.size = Pt(first_span["size"])
                        # tf.paragraphs[0].font.name = first_span["font"] # Font mapping is hard
                    except:
                        pass
                        
                elif block["type"] == 1: # Image
                    # Extract image
                    image_bytes = block["image"]
                    ext = block["ext"]
                    temp_img_name = f"temp_img_{i}_{int(b_x0)}_{int(b_y0)}.{ext}"
                    try:
                        with open(temp_img_name, "wb") as f:
                            f.write(image_bytes)
                        
                        slide.shapes.add_picture(temp_img_name, Pt(b_x0), Pt(b_y0), width=Pt(b_width), height=Pt(b_height))
                        os.remove(temp_img_name)
                    except Exception as e:
                        print(f"Warning: Failed to add image element: {e}")

        else:
            # Render page to image (pixmap)
            # Default matrix is 72 dpi, 2x zoom for better quality
            zoom = 2
            mat = fitz.Matrix(zoom, zoom)
            pix = page.get_pixmap(matrix=mat)
            
            # Save image to a temporary file
            temp_image_path = f"temp_page_{i}.png"
            pix.save(temp_image_path)
            
            # Get image dimensions from pixmap
            img_width_px = pix.width
            img_height_px = pix.height
            img_ratio = img_width_px / img_height_px
            
            slide_ratio = slide_width / slide_height
            
            # Calculate size and position to fit image within slide while preserving aspect ratio
            if img_ratio > slide_ratio:
                new_width = slide_width
                new_height = int(slide_width / img_ratio)
                top = int((slide_height - new_height) / 2)
                left = 0
                pic = slide.shapes.add_picture(temp_image_path, left, top, width=new_width)
            else:
                new_height = slide_height
                new_width = int(slide_height * img_ratio)
                left = int((slide_width - new_width) / 2)
                top = 0
                pic = slide.shapes.add_picture(temp_image_path, left, top, height=new_height)

            # Cleanup temp file
            os.remove(temp_image_path)

    try:
        prs.save(pptx_path)
        print(f"Successfully saved PPTX to {pptx_path}")
    except Exception as e:
        print(f"Error saving PPTX file: {e}")

if __name__ == "__main__":
    parser = argparse.ArgumentParser(description="Convert PDF to PPTX")
    parser.add_argument("pdf_path", help="Path to input PDF file")
    parser.add_argument("pptx_path", help="Path to output PPTX file")
    parser.add_argument("--editable", action="store_true", help="Attempt to reconstruct text as editable elements")
    
    args = parser.parse_args()
    
    convert_pdf_to_pptx(args.pdf_path, args.pptx_path, args.editable)
